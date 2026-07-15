"""
Thesis defence slide-deck generator.

Builds "Hybrid Static-Trace Analysis for Interaction Model Synthesis in
Event-Driven Systems" defence deck (16:9) using python-pptx.

Features:
  - Real TNO-ESI + UvA MSc-SE logos on title / closing slides
  - Dense H1 result table split across two clean slides
  - Speaker notes (talking points + timing) on every slide

Run:  python build_defence_slides.py
Output: Thesis_Defence_Prathik.pptx  (in the same folder)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as _RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def RGBColor(*a):
    """Accept either RGBColor(r, g, b) or RGBColor(0xRRGGBB)."""
    if len(a) == 1:
        v = int(a[0])
        return _RGBColor((v >> 16) & 0xFF, (v >> 8) & 0xFF, v & 0xFF)
    return _RGBColor(*a)

# --------------------------------------------------------------------------- #
#  Theme
# --------------------------------------------------------------------------- #
HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
LOGO_TNO = os.path.join(FIG, "tno_esi.jpg")
LOGO_MSC = os.path.join(FIG, "msc-se-logo-small-white.png")

NAVY   = RGBColor(0x0B, 0x2B, 0x45)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
TEAL   = RGBColor(0x1B, 0xA1, 0x9C)
AMBER  = RGBColor(0xE1, 0x8A, 0x1B)
GREY   = RGBColor(0x5A, 0x63, 0x6E)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
RED    = RGBColor(0xC0, 0x39, 0x2B)
DARK   = RGBColor(0x22, 0x2A, 0x33)

FONT   = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# running page counter so slides renumber automatically
_PG = 0
def pg():
    global _PG
    _PG += 1
    return _PG


# --------------------------------------------------------------------------- #
#  Helpers
# --------------------------------------------------------------------------- #
def slide():
    return prs.slides.add_slide(BLANK)


def note(s, text):
    """Attach speaker notes to a slide."""
    s.notes_slide.notes_text_frame.text = text.strip()


def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp


def textbox(s, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        p.level = ln.get("level", 0)
        txt = ln.get("text", "")
        if ln.get("bullet"):
            txt = "•  " + txt
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = ln.get("font", FONT)
        f.size = Pt(ln.get("size", 18))
        f.bold = ln.get("bold", False)
        f.italic = ln.get("italic", False)
        f.color.rgb = ln.get("color", DARK)
    return tb


def title_bar(s, title, kicker=None, num=None):
    rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    rect(s, 0, Inches(1.15), SW, Inches(0.06), fill=AMBER)
    if kicker:
        textbox(s, Inches(0.55), Inches(0.14), Inches(11), Inches(0.3),
                [{"text": kicker.upper(), "size": 12, "bold": True, "color": TEAL}])
    textbox(s, Inches(0.55), Inches(0.40), Inches(11.2), Inches(0.7),
            [{"text": title, "size": 27, "bold": True, "color": WHITE,
              "font": FONT_H}])
    if num is not None:
        textbox(s, Inches(12.2), Inches(0.42), Inches(0.9), Inches(0.5),
                [{"text": str(num), "size": 15, "bold": True, "color": TEAL,
                  "align": PP_ALIGN.RIGHT}])


def footer(s, num):
    textbox(s, Inches(0.4), Inches(7.05), Inches(9), Inches(0.35),
            [{"text": "Hybrid Static–Trace Interaction Model Synthesis  ·  Prathik Anand Krishnan  ·  TNO-ESI / UvA",
              "size": 9, "color": GREY}])
    textbox(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
            [{"text": str(num), "size": 10, "color": GREY, "align": PP_ALIGN.RIGHT}])


def panel(s, x, y, w, h, fill=LIGHT, line=None):
    r = rect(s, x, y, w, h, fill=fill,
             line=line, line_w=Pt(1) if line else None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        r.adjustments[0] = 0.04
    except Exception:
        pass
    return r


def chip(s, x, y, w, text, fill, tcolor=WHITE, h=Inches(0.42), size=13):
    r = rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        r.adjustments[0] = 0.25
    except Exception:
        pass
    tf = r.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = tcolor
    run.font.name = FONT
    return r


def img_fit(s, path, x, y, w, h):
    """Add picture constrained by the (w,h) box; keeps aspect ratio."""
    if not os.path.exists(path):
        rect(s, x, y, w, h, fill=LIGHT, line=GREY, line_w=Pt(1))
        textbox(s, x, y + h/2 - Inches(0.2), w, Inches(0.4),
                [{"text": "[figure: %s]" % os.path.basename(path),
                  "size": 12, "color": GREY, "align": PP_ALIGN.CENTER}])
        return
    pic = s.shapes.add_picture(path, x, y, width=w)
    if pic.height > h:
        scale = h / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = int(x + (w - pic.width) / 2)
    pic.top = int(y + (h - pic.height) / 2)
    return pic


IMG = img_fit


def logo(s, path, x, y, h, on_panel=False, panel_pad=Inches(0.12)):
    """Place a logo constrained by height; optionally on a white panel."""
    if not os.path.exists(path):
        return None
    if on_panel:
        # temp add to learn width
        tmp = s.shapes.add_picture(path, x, y, height=h)
        w = tmp.width
        tmp._element.getparent().remove(tmp._element)
        pw = w + 2 * panel_pad
        ph = h + 2 * panel_pad
        panel(s, x, y, pw, ph, fill=WHITE)
        pic = s.shapes.add_picture(path, x + panel_pad, y + panel_pad, height=h)
        return pic
    return s.shapes.add_picture(path, x, y, height=h)


def logo_right(s, path, right_x, y, h, on_panel=False):
    """Place a height-constrained logo right-aligned to right_x."""
    if not os.path.exists(path):
        return None
    tmp = s.shapes.add_picture(path, Inches(0), y, height=h)
    w = tmp.width
    tmp._element.getparent().remove(tmp._element)
    x = right_x - (w + (Inches(0.24) if on_panel else 0))
    return logo(s, path, x, y, h, on_panel=on_panel)


def table(s, x, y, w, rows, col_w, header_fill=NAVY, header_color=WHITE,
          row_h=Inches(0.42), font_size=13, highlight_rows=None,
          highlight_fill=RGBColor(0xEC, 0xF6, 0xEE)):
    highlight_rows = highlight_rows or []
    nrows = len(rows)
    ncols = len(rows[0])
    tot = sum(col_w)
    widths = [int(w * c / tot) for c in col_w]
    gt = s.shapes.add_table(nrows, ncols, x, y, w, row_h * nrows).table
    gt.first_row = False
    gt.horz_banding = False
    for ci, cw in enumerate(widths):
        gt.columns[ci].width = cw
    for ri in range(nrows):
        gt.rows[ri].height = row_h
        for ci in range(ncols):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            elif ri in highlight_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = highlight_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xF1, 0xF4, 0xF7)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(rows[ri][ci])
            r.font.size = Pt(font_size)
            r.font.name = FONT
            r.font.bold = (ri == 0) or (ri in highlight_rows and ci == 0)
            r.font.color.rgb = header_color if ri == 0 else DARK
    return gt


# --------------------------------------------------------------------------- #
#  SLIDE 1 — Title
# --------------------------------------------------------------------------- #
n = pg(); s = slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, Inches(0.35), SH, fill=AMBER)
rect(s, 0, Inches(5.05), SW, Inches(0.05), fill=TEAL)
# logos top-right
logo_right(s, LOGO_TNO, SW - Inches(0.6), Inches(0.55), Inches(0.85), on_panel=True)
logo(s, LOGO_MSC, Inches(0.9), Inches(6.55), Inches(0.6))
textbox(s, Inches(0.9), Inches(0.7), Inches(9.6), Inches(0.4),
        [{"text": "MSc Software Engineering  ·  Thesis Defence", "size": 15,
          "bold": True, "color": TEAL}])
textbox(s, Inches(0.9), Inches(1.5), Inches(11.6), Inches(2.2),
        [{"text": "Hybrid Static–Trace Analysis for", "size": 40, "bold": True,
          "color": WHITE, "font": FONT_H, "space_after": 2},
         {"text": "Interaction Model Synthesis in Event-Driven Systems",
          "size": 40, "bold": True, "color": WHITE, "font": FONT_H}])
textbox(s, Inches(0.9), Inches(3.55), Inches(11.6), Inches(0.6),
        [{"text": "Scenario-Specific Interaction Diagrams in SysML v2",
          "size": 20, "italic": True, "color": RGBColor(0xC9, 0xD6, 0xE2)}])
textbox(s, Inches(0.9), Inches(5.25), Inches(11.6), Inches(1.2),
        [{"text": "Prathik Anand Krishnan", "size": 22, "bold": True,
          "color": WHITE, "space_after": 6},
         {"text": "Examiner: Dr. L. Thomas van Binsbergen (CCI, UvA)   ·   "
                  "Reader: Prof. Dr. Benny Åkesson (PCS, UvA)",
          "size": 13, "color": RGBColor(0xC9, 0xD6, 0xE2), "space_after": 2},
         {"text": "Daily supervisor: Dr. Rosilde Corvino (TNO-ESI)",
          "size": 13, "color": RGBColor(0xC9, 0xD6, 0xE2)}])
textbox(s, Inches(2.0), Inches(6.62), Inches(7.5), Inches(0.5),
        [{"text": "University of Amsterdam  ·  TNO-ESI (in partnership with Philips), Eindhoven",
          "size": 12, "bold": True, "color": TEAL}])
note(s, """
[0:30] Good [morning/afternoon] everyone. My name is Prathik Anand Krishnan,
and today I'll defend my MSc Software Engineering thesis, done at TNO-ESI and
the University of Amsterdam.
The title is "Hybrid Static–Trace Analysis for Interaction Model Synthesis in
Event-Driven Systems" — in short, automatically reconstructing scenario-specific
SysML v2 interaction diagrams from C++ code by combining two kinds of evidence.
Thanks to my committee: my examiner Dr. van Binsbergen, my reader Prof. Åkesson,
and my daily supervisor Dr. Corvino.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 2 — Agenda
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Agenda", kicker="Overview", num=n)
items = [
    ("1", "Motivation", "Why interaction models, why they drift"),
    ("2", "Problem & Research Questions", "The bidirectional fusion gap"),
    ("3", "Approach", "The GDI hybrid pipeline & property graph"),
    ("4", "Case Study", "CPSCore — event-driven C++ framework"),
    ("5", "Experiments", "H1 accuracy · H2 focus · H3 utility"),
    ("6", "Results & Findings", "When fusion helps — and when it doesn't"),
    ("7", "Contributions, Limitations & Conclusion", ""),
]
y = Inches(1.55)
for nn, t, d in items:
    chip(s, Inches(0.7), y, Inches(0.55), nn, BLUE, h=Inches(0.55), size=18)
    textbox(s, Inches(1.45), y - Inches(0.02), Inches(3.3), Inches(0.6),
            [{"text": t, "size": 18, "bold": True, "color": NAVY}])
    if d:
        textbox(s, Inches(4.9), y + Inches(0.06), Inches(8), Inches(0.4),
                [{"text": d, "size": 13, "color": GREY}])
    y += Inches(0.72)
footer(s, n)
note(s, """
[0:30] Here's the roadmap. I'll start with the motivation and the concrete
problem, state my two research questions, then walk through the approach —
the pipeline and the shared property graph. After a quick look at the case
study, CPSCore, I'll present the three experiments and their results, and close
with contributions, limitations, and conclusions.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 3 — Motivation
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Understanding Behaviour Before You Change It",
                       kicker="Motivation", num=n)
panel(s, Inches(0.55), Inches(1.5), Inches(6.05), Inches(4.9))
textbox(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(4.6),
        [{"text": "The setting", "size": 16, "bold": True, "color": BLUE,
          "space_after": 6},
         {"text": "Long-lived, safety-critical systems (medical imaging, "
          "industrial automation) evolve continuously.", "size": 15,
          "bullet": True, "color": DARK, "space_after": 8},
         {"text": "Behaviour is meant to be captured by Interaction Models — "
          "scenario-specific sequence diagrams.", "size": 15, "bullet": True,
          "space_after": 8},
         {"text": "In practice these diagrams drift out of sync with the code "
          "and become unreliable.", "size": 15, "bullet": True,
          "space_after": 8},
         {"text": "Engineers must reconstruct behaviour by hand — slow, "
          "error-prone, expertise-dependent.", "size": 15, "bullet": True,
          "space_after": 8},
         {"text": "Event-driven systems make this worse: async callbacks and "
          "pub/sub hide who talks to whom.", "size": 15, "bullet": True}])
panel(s, Inches(6.8), Inches(1.5), Inches(5.95), Inches(4.9), fill=RGBColor(0xFB,0xF2,0xE3))
textbox(s, Inches(7.1), Inches(1.72), Inches(5.4), Inches(0.5),
        [{"text": "A REAL DEFECT THIS WOULD HAVE CAUGHT", "size": 12,
          "bold": True, "color": AMBER}])
textbox(s, Inches(7.1), Inches(2.2), Inches(5.4), Inches(4.0),
        [{"text": "A stale proxy object was never freed — a reset call was "
          "missing from the connection-handling code.", "size": 15,
          "bullet": True, "color": DARK, "space_after": 10},
         {"text": "After a service restart, the system silently stopped "
          "delivering change-notification events.", "size": 15, "bullet": True,
          "space_after": 10},
         {"text": "The fault was visible only by following an implicit chain of "
          "events that no interaction model showed.", "size": 15,
          "bullet": True, "space_after": 10},
         {"text": "Diagnosing it meant manually reconstructing exactly the kind "
          "of model this thesis automates.", "size": 15, "bullet": True,
          "bold": True, "color": NAVY}])
footer(s, n)
note(s, """
[2:00] The systems I care about are long-lived and safety-critical — think
medical imaging or industrial automation. Their behaviour is supposed to be
captured in interaction models, i.e. sequence diagrams. But in practice those
diagrams drift out of sync as the code evolves, so engineers end up
reconstructing behaviour by hand — slow and expertise-dependent. Event-driven
systems make it harder because async callbacks and pub/sub hide who talks to
whom.
On the right is a real, previously-resolved defect that motivated the work: a
stale proxy was never freed because a reset call was missing, and after a
restart the system silently stopped delivering notifications. It was only
visible by tracing an implicit event chain that no diagram showed — exactly the
kind of model this thesis reconstructs automatically.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 4 — Problem statement
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Two Incomplete Sources of Evidence",
                       kicker="Problem Statement", num=n)
panel(s, Inches(0.55), Inches(1.55), Inches(5.9), Inches(2.35), fill=RGBColor(0xEAF1F8))
textbox(s, Inches(0.8), Inches(1.7), Inches(5.4), Inches(0.4),
        [{"text": "STATIC CALL GRAPH", "size": 14, "bold": True, "color": BLUE}])
textbox(s, Inches(0.8), Inches(2.15), Inches(5.4), Inches(1.7),
        [{"text": "Complete view of all possible interactions.", "size": 14,
          "bullet": True, "color": GREEN, "space_after": 5},
         {"text": "Over-approximates — includes paths never taken.", "size": 14,
          "bullet": True, "color": RED, "space_after": 5},
         {"text": "No execution order, no scenario specificity.", "size": 14,
          "bullet": True, "color": RED}])
panel(s, Inches(6.85), Inches(1.55), Inches(5.9), Inches(2.35), fill=RGBColor(0xEAF6F1))
textbox(s, Inches(7.1), Inches(1.7), Inches(5.4), Inches(0.4),
        [{"text": "RUNTIME TRACES", "size": 14, "bold": True, "color": TEAL}])
textbox(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(1.7),
        [{"text": "Scenario-specific by construction, real order.", "size": 14,
          "bullet": True, "color": GREEN, "space_after": 5},
         {"text": "Under-approximates — only what was exercised.", "size": 14,
          "bullet": True, "color": RED, "space_after": 5},
         {"text": "Sparse / noisy; misses suppressed or rare paths.", "size": 14,
          "bullet": True, "color": RED}])
panel(s, Inches(0.55), Inches(4.1), Inches(12.2), Inches(2.3), fill=NAVY)
textbox(s, Inches(0.85), Inches(4.28), Inches(11.6), Inches(0.5),
        [{"text": "THE GAP", "size": 13, "bold": True, "color": AMBER}])
textbox(s, Inches(0.85), Inches(4.7), Inches(11.6), Inches(1.6),
        [{"text": "Existing hybrid methods use static structure only to organise "
          "an already-observed trace.", "size": 16, "color": WHITE,
          "bullet": True, "space_after": 8},
         {"text": "No approach combines both directions at once: use static "
          "structure to recover interactions the trace missed, AND use the "
          "trace to trim the over-approximate static graph to the scenario.",
          "size": 16, "color": WHITE, "bullet": True, "bold": True},
         {"text": "Plus: publish/subscribe indirection (boost::signals2, Redis) "
          "hides the subscription topology from call-graph analysis entirely.",
          "size": 15, "color": RGBColor(0xC9,0xD6,0xE2), "bullet": True,
          "space_before": 6}])
footer(s, n)
note(s, """
[2:00] Two classic sources of evidence, each incomplete. The static call graph
is complete in the sense that it captures every possible interaction — but it
over-approximates and has no execution order or scenario specificity. Runtime
traces are the opposite: scenario-specific and correctly ordered, but
under-approximate — they only show what was actually exercised, and are sparse
or noisy.
The gap, in the box: prior hybrid methods only use static structure to organise
a trace they already have. Nobody combines both directions at once — using
static structure to recover what the trace missed, AND using the trace to trim
the static graph down to the scenario. And pub/sub indirection hides the
subscription topology from call-graph analysis entirely. Closing that
two-directional gap is the core problem.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 5 — Research questions
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Research Questions & Hypotheses",
                       kicker="Objectives", num=n)
data = [
    ("RQ1", "How does combining static analysis and incomplete runtime traces "
            "improve the accuracy and completeness of scenario-specific "
            "interaction model reconstruction?", "H1"),
    ("RQ2", "How can scenario-scoped querying of the static call graph "
            "restrict reconstruction to only the behaviour relevant to a "
            "given scenario?", "H2"),
]
y = Inches(1.5)
for rq, q, ans in data:
    panel(s, Inches(0.55), y, Inches(12.2), Inches(1.15))
    chip(s, Inches(0.75), y + Inches(0.3), Inches(1.15), rq, BLUE,
         h=Inches(0.55), size=18)
    textbox(s, Inches(2.1), y + Inches(0.14), Inches(8.1), Inches(0.95),
            [{"text": q, "size": 15, "bold": True, "color": NAVY}],
            anchor=MSO_ANCHOR.MIDDLE)
    chip(s, Inches(10.4), y + Inches(0.34), Inches(2.15), ans, TEAL,
         h=Inches(0.48), size=13)
    y += Inches(1.32)
panel(s, Inches(0.55), y, Inches(12.2), Inches(0.9), fill=RGBColor(0xFB,0xF2,0xE3))
textbox(s, Inches(0.8), y + Inches(0.12), Inches(12), Inches(0.7),
        [{"text": "H3 (downstream): the reconstructed diagram improves — and can "
          "even be sufficient for — LLM-agent scenario comprehension.",
          "size": 15, "bold": True, "color": AMBER}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, n)
note(s, """
[1:30] Two research questions. RQ1 is about accuracy: does combining static
analysis with incomplete traces actually improve reconstruction accuracy and
completeness? RQ2 is about focus: can scenario-scoped querying of the static
call graph restrict reconstruction to only scenario-relevant behaviour?
Each RQ has a matching hypothesis and a dedicated experiment — H1 answers RQ1,
H2 answers RQ2. And a downstream question, H3: is the reconstructed diagram
actually useful to a consumer — here, an LLM agent doing scenario comprehension.
""")

# --------------------------------------------------------------------------- #
#  SLIDE — Related Work & Research Gap  (from proposal slide 4)
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "State of the Art & the Research Gap",
                       kicker="Related Work", num=n)
cols = [
    ("Trace-Based", TEAL,
     [("+", "Concrete, observed behaviour"),
      ("x", "Traces are incomplete"),
      ("x", "Misses unexercised interactions")]),
    ("Static Analysis", BLUE,
     [("+", "Complete static structure"),
      ("x", "Over-approximates behaviour"),
      ("x", "No scenario focus")]),
    ("Existing Hybrid", GREY,
     [("+", "Uses static to reduce logging"),
      ("x", "Assumes observed traces complete"),
      ("x", "Does not recover missing edges")]),
    ("This Thesis", AMBER,
     [("+", "Recovers trace-missed edges from static"),
      ("+", "Trims static to scenario via trace"),
      ("+", "Recovers pub/sub wiring")]),
]
cw = Inches(2.98); gapx = Inches(0.13); x = Inches(0.55); y = Inches(1.5); ch = Inches(3.5)
for name, col, pts in cols:
    highlight = (name == "This Thesis")
    panel(s, x, y, cw, ch, fill=RGBColor(0xFB,0xF2,0xE3) if highlight else LIGHT,
          line=AMBER if highlight else None)
    rect(s, x, y, cw, Inches(0.6), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x, y + Inches(0.08), cw, Inches(0.45),
            [{"text": name, "size": 14, "bold": True,
              "color": NAVY if highlight else WHITE, "align": PP_ALIGN.CENTER}])
    yy = y + Inches(0.8)
    for mark, txt in pts:
        good = (mark == "+")
        textbox(s, x + Inches(0.12), yy, Inches(0.3), Inches(0.8),
                [{"text": "\u2714" if good else "\u2716", "size": 15, "bold": True,
                  "color": GREEN if good else RED}])
        textbox(s, x + Inches(0.42), yy, cw - Inches(0.55), Inches(0.85),
                [{"text": txt, "size": 12.5, "color": DARK}])
        yy += Inches(0.85)
    x += cw + gapx
panel(s, Inches(0.55), Inches(5.25), Inches(12.2), Inches(1.65), fill=NAVY)
textbox(s, Inches(0.85), Inches(5.42), Inches(11.7), Inches(0.4),
        [{"text": "THE GAP", "size": 13, "bold": True, "color": AMBER}])
textbox(s, Inches(0.85), Inches(5.82), Inches(11.7), Inches(1.0),
        [{"text": "No existing approach infers unobserved interactions by "
          "combining static structure and runtime traces to generate "
          "scenario-specific interaction models — closing this gap in both "
          "directions is the core contribution.", "size": 16, "bold": True,
          "color": WHITE}])
footer(s, n)
note(s, """
[1:30] Positioning the work. Three families exist today. Trace-based methods give
concrete observed behaviour but are incomplete and miss anything not exercised.
Static analysis gives complete structure but over-approximates and has no
scenario focus. Existing hybrids are the closest, but they use static structure
only to reduce logging or organise a trace — they assume the observed trace is
complete and never recover the interactions it missed.
The fourth column is what this thesis adds: it recovers trace-missed edges from
static evidence, trims the over-approximate static graph to the scenario using
the trace, and recovers pub/sub wiring invisible to call graphs. The gap, stated
plainly at the bottom: no prior approach infers unobserved interactions by
combining both sources for scenario-specific models — and closing it in both
directions is the core contribution.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 6 — Pipeline
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "The GDI Hybrid Pipeline", kicker="Approach", num=n)
stages = [
    ("Static\nExtraction", "Renaissance\nC++ → graph", BLUE),
    ("Property\nGraph", "Neo4j\n11 / 11 schema", NAVY),
    ("Runtime\nTracing", "csp_matcher\nprobes", TEAL),
    ("Graph\nFusion", "static ∪ trace\n(provenance)", BLUE),
    ("Scenario\nScoping", "scope-constrained\nCypher", NAVY),
    ("Agent\nSynthesis", "provenance-aware\npruning", TEAL),
    ("SysML v2 /\nMermaid", "scenario\ndiagram", AMBER),
]
ns = len(stages)
gap = Inches(0.2)
total_w = Inches(12.4)
bw = int((total_w - gap * (ns - 1)) / ns)
x = Inches(0.5); y = Inches(2.35); bh = Inches(1.7)
for i, (t, sub, c) in enumerate(stages):
    panel(s, x, y, bw, bh, fill=c)
    textbox(s, x, y + Inches(0.22), bw, Inches(0.9),
            [{"text": t, "size": 14, "bold": True, "color": WHITE,
              "align": PP_ALIGN.CENTER}])
    textbox(s, x, y + Inches(1.02), bw, Inches(0.6),
            [{"text": sub, "size": 10.5, "color": RGBColor(0xDCE6F0),
              "align": PP_ALIGN.CENTER}])
    if i < ns - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                x + bw + Emu(int(gap*0.05)),
                                y + bh/2 - Inches(0.12),
                                gap - Emu(int(gap*0.1)), Inches(0.24))
        ar.fill.solid(); ar.fill.fore_color.rgb = GREY
        ar.line.fill.background(); ar.shadow.inherit = False
    x += bw + gap
textbox(s, Inches(0.5), Inches(1.5), Inches(11), Inches(0.6),
        [{"text": "11 composable agent skills — each invocable by natural "
          "language, no code changes to the target system.", "size": 15,
          "bold": True, "color": NAVY}])
textbox(s, Inches(0.5), Inches(1.92), Inches(12.3), Inches(0.35),
        [{"text": "Originally scoped as a 5-stage methodology (proposal); "
          "realised as this graph-centric 7-stage pipeline.", "size": 12,
          "italic": True, "color": GREY}])
panel(s, Inches(0.5), Inches(4.5), Inches(6.05), Inches(1.7), fill=LIGHT)
textbox(s, Inches(0.75), Inches(4.65), Inches(5.6), Inches(1.5),
        [{"text": "Fusion (RQ1)", "size": 14, "bold": True, "color": BLUE,
          "space_after": 3},
         {"text": "Set union of static (CppCalls) and trace (TRACE_CALLS) edges, "
          "each tagged with which source observed it.", "size": 13.5,
          "color": DARK}])
panel(s, Inches(6.7), Inches(4.5), Inches(6.05), Inches(1.7), fill=LIGHT)
textbox(s, Inches(6.95), Inches(4.65), Inches(5.6), Inches(1.5),
        [{"text": "Scoping (RQ2)", "size": 14, "bold": True, "color": TEAL,
          "space_after": 3},
         {"text": "Scenario-scoped Cypher query: keep e where src ∈ "
          "scenario.primary and tgt ∈ scenario.targets — identical over "
          "static & trace edges.", "size": 13.5, "color": DARK}])
footer(s, n)
note(s, """
[2:30] This is the whole pipeline. It's implemented as eleven composable agent
skills, each invocable in natural language, and — importantly — with no changes
to the target system's own code.
Left to right: static extraction with the Renaissance tool turns C++ into a
graph; everything lands in a Neo4j property graph; csp_matcher inserts runtime
trace probes; graph fusion takes the union of static and trace edges, tagging
each with its provenance; scenario scoping filters that fused graph; an agent
synthesis step prunes false positives using provenance; and finally we emit a
SysML v2 or Mermaid diagram.
The two conceptual cores are at the bottom: fusion answers RQ1, scoping answers
RQ2. Scoping is just a scenario-scoped Cypher query that works identically over
static and trace edges — I'll show why on the next slide.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 7 — Property graph
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "A Single Property Graph for Both Evidence Types",
                       kicker="Approach · Foundations", num=n)
textbox(s, Inches(0.55), Inches(1.35), Inches(12), Inches(0.5),
        [{"text": "Static call graph and runtime trace live in one Neo4j graph — "
          "each interaction remembers which evidence produced it.",
          "size": 16, "bold": True, "color": NAVY}])
mets = [("11", "node labels"), ("11", "relationship types"),
        ("2,574", "ordered static interactions"),
        ("121", "components"),
        ("27", "runtime trace events"),
        ("26", "unique traced pairs")]
y = Inches(2.05)
for v, lbl in mets:
    panel(s, Inches(0.55), y, Inches(5.4), Inches(0.62), fill=LIGHT)
    textbox(s, Inches(0.7), y + Inches(0.05), Inches(1.7), Inches(0.5),
            [{"text": v, "size": 20, "bold": True, "color": BLUE,
              "align": PP_ALIGN.RIGHT}])
    textbox(s, Inches(2.5), y + Inches(0.12), Inches(3.4), Inches(0.4),
            [{"text": lbl, "size": 14, "color": DARK}])
    y += Inches(0.72)
panel(s, Inches(6.3), Inches(2.05), Inches(6.45), Inches(4.35), fill=NAVY)
textbox(s, Inches(6.6), Inches(2.25), Inches(5.9), Inches(0.4),
        [{"text": "FUSION & SCOPE, FORMALLY", "size": 13, "bold": True,
          "color": AMBER}])
textbox(s, Inches(6.6), Inches(2.75), Inches(5.9), Inches(1.2),
        [{"text": "Gσ = scope(G_static ∪ G_trace, σ)", "size": 17, "bold": True,
          "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 6},
         {"text": "= { e ∈ Es ∪ Et  |  src(e) ∈ σ.primary  ∧  tgt(e) ∈ σ.targets }",
          "size": 14, "color": RGBColor(0xC9,0xD6,0xE2),
          "align": PP_ALIGN.CENTER}])
textbox(s, Inches(6.6), Inches(4.15), Inches(5.9), Inches(2.1),
        [{"text": "Both edge types carry ClientComponent / ServerComponent "
          "metadata, so one filter works over both.", "size": 14,
          "color": WHITE, "bullet": True, "space_after": 8},
         {"text": "Schema-introspection skill reads labels at runtime → the "
          "pipeline is codebase-agnostic.", "size": 14, "color": WHITE,
          "bullet": True, "space_after": 8},
         {"text": "Recovers pub/sub wiring (boost::signals2) that plain "
          "call-graph analysis cannot see.", "size": 14, "color": WHITE,
          "bullet": True}])
footer(s, n)
note(s, """
[2:00] This underlies both RQ1 and RQ2. The key design decision is that both
evidence types live in one Neo4j property graph, and every interaction edge
remembers which source produced it — its provenance. The schema has 11 node
labels and 11 relationship types. On CPSCore the static side yields 2,574
ordered interactions across 121 components; the trace side adds 27 events over
26 unique pairs.
On the right is the formal operation: the scenario graph is the scenario-scoped
subset of the union of static and trace edges, keeping edges whose source is in
the scenario's primary component and whose target is a declared participant.
Because both edge types carry client/server component metadata, one linear
filter works over both.
A schema-introspection skill reads the labels at runtime, so the pipeline is
codebase-agnostic — and crucially, the graph lets us recover pub/sub wiring that
plain call-graph analysis can't see.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 8 — Static extraction
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Static Extraction: The Full Call Graph",
                       kicker="Approach · Static", num=n)
IMG(s, os.path.join(FIG, "function-call-dependencies.png"),
    Inches(0.55), Inches(1.4), Inches(8.4), Inches(5.3))
panel(s, Inches(9.2), Inches(1.5), Inches(3.55), Inches(5.0), fill=LIGHT)
textbox(s, Inches(9.45), Inches(1.7), Inches(3.1), Inches(4.7),
        [{"text": "What it gives us", "size": 15, "bold": True, "color": BLUE,
          "space_after": 8},
         {"text": "Renaissance parses C++ → Neo4j.", "size": 13.5,
          "bullet": True, "space_after": 7},
         {"text": "Over-approximate: every possible call, scenario-agnostic.",
          "size": 13.5, "bullet": True, "space_after": 7},
         {"text": "Projected into an ordered sequence-dependency table "
          "(2,574 interactions).", "size": 13.5, "bullet": True,
          "space_after": 7},
         {"text": "Recall floor: guarantees interactions a sparse trace would "
          "miss.", "size": 13.5, "bullet": True, "bold": True,
          "color": NAVY}])
footer(s, n)
note(s, """
[1:30] The static half. The Renaissance semantic-extraction tool parses the C++
and populates Neo4j. What you're seeing is the full function-call dependency
graph — everything the code could do, scenario-agnostic, so it over-approximates.
We project the CppCalls edges into an ordered sequence-dependency table — 2,574
interactions. The important role of static analysis in the fusion is as a recall
floor: it guarantees interactions that a sparse trace would otherwise miss.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 9 — Runtime tracing
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Runtime Tracing with csp_matcher",
                       kicker="Approach · Dynamic", num=n)
IMG(s, os.path.join(FIG, "signal-runtime-tracing.png"),
    Inches(6.5), Inches(1.45), Inches(6.3), Inches(5.1))
textbox(s, Inches(0.55), Inches(1.45), Inches(5.7), Inches(5.0),
        [{"text": "A small, reversible instrumentation tool", "size": 16,
          "bold": True, "color": NAVY, "space_after": 10},
         {"text": "Clang LibTooling-based; inserts trace probes by AST-level "
          "pattern matching.", "size": 15, "bullet": True, "space_after": 9},
         {"text": "Guided by the code graph — probes only the call sites that "
          "matter for a scenario.", "size": 15, "bullet": True,
          "space_after": 9},
         {"text": "Fully reversible: add_tracing / remove_tracing restore the "
          "original source exactly.", "size": 15, "bullet": True,
          "space_after": 9},
         {"text": "Recovers pub/sub edges: resolves boost::signals2 "
          "connect()/slot wiring invisible to static call sites.", "size": 15,
          "bullet": True, "space_after": 9},
         {"text": "Produced 27 trace events across the studied scenarios.",
          "size": 15, "bullet": True, "bold": True, "color": BLUE}])
footer(s, n)
note(s, """
[2:00] The dynamic half is captured by csp_matcher, a Clang LibTooling tool I
built. It inserts trace probes by matching AST patterns, and it's guided by the
code graph so it only instruments the call sites relevant to a scenario. It's
fully reversible — add and remove leave the source byte-for-byte unchanged.
The figure illustrates the pub/sub case: csp_matcher resolves boost::signals2
connect and slot wiring, which has no static call site, so it recovers edges the
call graph simply cannot see. On CPSCore it produced 27 trace events.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 10 — CPSCore
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Case Study: CPSCore", kicker="Evaluation Target",
                       num=n)
IMG(s, os.path.join(FIG, "module-dependency.png"),
    Inches(6.5), Inches(1.4), Inches(6.3), Inches(5.2))
textbox(s, Inches(0.55), Inches(1.4), Inches(5.7), Inches(2.2),
        [{"text": "Open-source event-driven C++ framework for cyber-physical "
          "systems.", "size": 15, "bold": True, "color": NAVY,
          "space_after": 8},
         {"text": "~13,000 LOC · 154 files · 5 modules: Aggregation, "
          "Configuration, Synchronization, Logging, Utilities.", "size": 14,
          "bullet": True, "space_after": 6},
         {"text": "Event-driven pub/sub via boost::signals2 and Redis.",
          "size": 14, "bullet": True}])
panel(s, Inches(0.55), Inches(3.75), Inches(5.7), Inches(2.65),
      fill=RGBColor(0xFB,0xF2,0xE3))
textbox(s, Inches(0.8), Inches(3.9), Inches(5.2), Inches(2.4),
        [{"text": "KEY STRUCTURAL PROPERTY", "size": 12, "bold": True,
          "color": AMBER, "space_after": 6},
         {"text": "~470 template occurrences (static polymorphism); only 32 "
          "virtual signatures, all within one module or at config boundaries.",
          "size": 13.5, "bullet": True, "space_after": 6},
         {"text": "No cross-module call goes through a vtable → static graph is "
          "reachability-complete w.r.t. the trace on found-in-the-wild "
          "scenarios.", "size": 13.5, "bullet": True, "space_after": 6},
         {"text": "This is exactly why we constructed scenario S4.", "size": 13.5,
          "bold": True, "color": NAVY, "bullet": True}])
footer(s, n)
note(s, """
[1:30] The primary case study is CPSCore — an open-source event-driven C++
framework for cyber-physical systems: about 13,000 lines, 154 files, five
modules, with pub/sub over boost::signals2 and Redis.
One structural property matters a lot for interpreting the results, in the amber
box: CPSCore leans heavily on templates — static polymorphism — with only 32
virtual signatures, all inside a single module or at config boundaries. No
cross-module call goes through a vtable. That means, on the naturally-occurring
scenarios, the static graph is already reachability-complete with respect to the
trace. That's a property of this system, not a law — and it's precisely why I
constructed scenario S4 to break it.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 11 — Scenarios
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Four Evaluation Scenarios", kicker="Evaluation", num=n)
IMG(s, os.path.join(FIG, "s4-structural.png"),
    Inches(7.1), Inches(1.5), Inches(5.7), Inches(3.4))
rows = [
    ["Scenario", "Nature", "What it exercises"],
    ["S1", "Happy-path", "Aggregation chain"],
    ["S2", "Happy-path", "Config mapping"],
    ["S3", "Happy-path", "Multi-component runner"],
    ["S4", "Constructed", "pub/sub dual-gap fault path"],
]
table(s, Inches(0.55), Inches(1.55), Inches(6.2), rows, [1.1, 1.4, 2.7],
      row_h=Inches(0.55), font_size=13, highlight_rows=[4])
panel(s, Inches(0.55), Inches(4.9), Inches(6.2), Inches(1.55), fill=NAVY)
textbox(s, Inches(0.78), Inches(5.05), Inches(5.75), Inches(1.35),
        [{"text": "S4 — the dual-gap scenario", "size": 13, "bold": True,
          "color": AMBER, "space_after": 4},
         {"text": "A boost::signals2 bridge where static sees stream but not "
          "onStageEvent (runtime connect()), and dynamic sees onStageEvent but "
          "not stream (a LogLevel::NONE no-op). Only the union recovers both.",
          "size": 13, "color": WHITE}])
textbox(s, Inches(7.1), Inches(4.95), Inches(5.7), Inches(0.5),
        [{"text": "S4 structural view — StageEventBridge ↔ StageEventListener",
          "size": 11, "italic": True, "color": GREY, "align": PP_ALIGN.CENTER}])
footer(s, n)
note(s, """
[2:00] Four scenarios. S1 to S3 are happy-path scenarios that occur naturally in
CPSCore — an aggregation chain, a config mapping, and a multi-component runner.
S4 is deliberately constructed: a boost::signals2 publish/subscribe bridge
designed so each single source is blind to a *different* edge. Static sees the
`stream` call but not `onStageEvent`, because the slot is connected at runtime
via connect() with no static call expression. Dynamic sees `onStageEvent` but
not `stream`, because a LogLevel::NONE idiom makes the log call a runtime no-op.
Only the union recovers both. S4 isolates the one regime where fusion is
strictly additive — I report it as its own stratum because it's constructed, not
found in the wild.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 12 — H1 setup
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "H1 — Does Fusion Improve Accuracy?",
                       kicker="Experiment · RQ1", num=n)
textbox(s, Inches(0.55), Inches(1.35), Inches(12), Inches(0.5),
        [{"text": "Five conditions × four scenarios, scored against a manually "
          "source-read reference edge set (authored before scoring).",
          "size": 15, "bold": True, "color": NAVY}])
conds = [
    ("C1", "LLM-only", "no code context", GREY),
    ("C2", "Static-only", "call graph", BLUE),
    ("C3", "Dynamic-only", "runtime traces", TEAL),
    ("C4", "Static ∪ Dynamic", "raw union", NAVY),
    ("C5", "Full pipeline", "union + agent pruning", AMBER),
]
y = Inches(2.05)
for c, t, d, col in conds:
    panel(s, Inches(0.55), y, Inches(7.4), Inches(0.72), fill=LIGHT)
    chip(s, Inches(0.72), y + Inches(0.14), Inches(0.9), c, col, h=Inches(0.44),
         size=15)
    textbox(s, Inches(1.8), y + Inches(0.08), Inches(3.0), Inches(0.55),
            [{"text": t, "size": 15, "bold": True, "color": NAVY}],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.7), y + Inches(0.08), Inches(3.1), Inches(0.55),
            [{"text": d, "size": 13, "color": GREY}], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.82)
panel(s, Inches(8.25), Inches(2.05), Inches(4.5), Inches(4.15), fill=NAVY)
textbox(s, Inches(8.5), Inches(2.25), Inches(4.0), Inches(0.4),
        [{"text": "METRICS (on interaction edges)", "size": 12.5, "bold": True,
          "color": AMBER}])
textbox(s, Inches(8.5), Inches(2.75), Inches(4.0), Inches(3.3),
        [{"text": "P = |gen ∩ ref| / |gen|", "size": 15, "bold": True,
          "color": WHITE, "space_after": 8},
         {"text": "R = |gen ∩ ref| / |ref|", "size": 15, "bold": True,
          "color": WHITE, "space_after": 8},
         {"text": "F1 = 2PR / (P + R)", "size": 15, "bold": True, "color": WHITE,
          "space_after": 14},
         {"text": "Aggregated as macro-F1 — each scenario weighted equally.",
          "size": 13.5, "color": RGBColor(0xC9,0xD6,0xE2)},
         {"text": "S4 kept as its own stratum (constructed, not found in the "
          "wild).", "size": 13.5, "color": RGBColor(0xC9,0xD6,0xE2),
          "space_before": 8}])
footer(s, n)
note(s, """
[1:30] H1 tests RQ1. I compare five conditions across the four scenarios. C1 is
an LLM-only baseline with no code context. C2 is static-only, C3 is dynamic-only.
C4 is the raw union of the two. C5 is the full pipeline — the union plus a
provenance-aware agent step that prunes false positives.
Everything is scored on interaction edges — unique ordered sender/receiver pairs
— against a reference set I authored by reading the source *before* scoring
anything. Standard precision, recall, F1, aggregated as macro-F1 so each
scenario counts equally. S4 is reported as its own stratum.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 13 — H1 results (overall P/R/F1) — split part 1
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "H1 Results — Overall Precision / Recall / F1",
                       kicker="Results · RQ1  (S1–S4)", num=n)
rows = [
    ["Condition", "Precision", "Recall", "F1"],
    ["C1  LLM-only", "0.13", "0.13", "0.13"],
    ["C2  Static-only", "0.90", "0.88", "0.85"],
    ["C3  Dynamic-only", "1.00", "0.88", "0.92"],
    ["C4  Static ∪ Dynamic", "0.90", "1.00", "0.94"],
    ["C5  Full pipeline", "0.94", "1.00", "0.96"],
]
table(s, Inches(1.4), Inches(1.65), Inches(10.5), rows, [3.2, 2.0, 2.0, 2.0],
      row_h=Inches(0.66), font_size=16, highlight_rows=[5])
panel(s, Inches(1.4), Inches(5.9), Inches(10.5), Inches(1.0),
      fill=RGBColor(0xEC,0xF6,0xEE))
textbox(s, Inches(1.65), Inches(6.02), Inches(10.1), Inches(0.8),
        [{"text": "C5 posts the highest overall F1 of any condition (0.96): it "
          "keeps C4's perfect recall while partially fixing static "
          "over-approximation.", "size": 15, "bold": True, "color": GREEN}],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, n)
note(s, """
[2:00] First the headline, aggregated over all four scenarios, one number per
cell. Read down the F1 column: the LLM baseline is near-useless at 0.13. Static
alone is 0.85, dynamic alone 0.92. The raw union C4 lifts recall to a perfect
1.00 and F1 to 0.94. And the full pipeline C5 is the best condition tested at
0.96 — it preserves C4's perfect recall while partially correcting static
over-approximation. Note dynamic-only already has perfect precision, but its
recall is only 0.88 — it misses edges — which is exactly what fusion fixes. The
*where* and *why* are on the next slide.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 14 — H1 results per-stratum F1 — split part 2
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "H1 Results — Fusion Is Additive Only on S4",
                       kicker="Results · RQ1  (per stratum, F1)", num=n)
rows = [
    ["Condition", "Happy-path (S1–S3)", "Constructed (S4)", "Overall"],
    ["C1  LLM-only", "0.00", "0.50", "0.13"],
    ["C2  Static-only", "0.92", "0.67", "0.85"],
    ["C3  Dynamic-only", "1.00", "0.67", "0.92"],
    ["C4  Static ∪ Dynamic", "0.92", "1.00", "0.94"],
    ["C5  Full pipeline", "0.95", "1.00", "0.96"],
]
table(s, Inches(0.7), Inches(1.6), Inches(11.9), rows, [3.0, 3.0, 2.9, 2.0],
      row_h=Inches(0.6), font_size=15, highlight_rows=[4, 5])
panel(s, Inches(0.7), Inches(5.5), Inches(5.85), Inches(1.4), fill=NAVY)
textbox(s, Inches(0.95), Inches(5.62), Inches(5.4), Inches(1.2),
        [{"text": "On S4: C4 = 1.00 strictly beats C2 = 0.67 and C3 = 0.67.",
          "size": 14, "bold": True, "color": WHITE, "space_after": 4},
         {"text": "The only scenario where fusion is additive — each source is "
          "blind to a different edge.", "size": 12.5,
          "color": RGBColor(0xC9,0xD6,0xE2)}])
panel(s, Inches(6.75), Inches(5.5), Inches(5.85), Inches(1.4),
      fill=RGBColor(0xEAF1F8))
textbox(s, Inches(7.0), Inches(5.62), Inches(5.4), Inches(1.2),
        [{"text": "On S1–S3: C4 = C2 (fusion redundant).", "size": 14,
          "bold": True, "color": BLUE, "space_after": 4},
         {"text": "C3 ⊆ C2 here, so the union adds nothing — but it still "
          "guarantees a recall floor (C4 ≥ C3 always).", "size": 12.5,
          "color": DARK}])
footer(s, n)
note(s, """
[2:00] Now the same F1 split by stratum, which is the heart of H1. Look at the
S4 column: C4 scores a perfect 1.00, strictly above both C2 and C3 at 0.67. That
is the single scenario where fusion is genuinely additive — because static and
dynamic are each blind to a different edge, and only the union recovers both.
On the happy-path stratum, C4 equals C2. That's not a flaw: because there's no
cross-module virtual dispatch, the trace edges are a subset of the static edges,
so the union adds nothing new. But fusion is still protective — it guarantees a
recall floor whenever traces are sparse. The takeaway: fusion helps exactly
where the sources disagree, and is safely redundant otherwise. In any system
with virtual dispatch, plugin loading, or reflection, the additive regime is
what you'd see.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 15 — H1 interpretation
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Why the Numbers Look the Way They Do",
                       kicker="Results · Interpretation", num=n)
cards = [
    ("Fusion additive on S4", GREEN,
     "Static misses the runtime connect() edge; dynamic misses the "
     "LogLevel::NONE-suppressed edge. Only the union C4 recovers both → "
     "F1 1.00 vs 0.67 / 0.67."),
    ("Fusion redundant on S1–S3", BLUE,
     "No cross-module virtual dispatch in CPSCore ⇒ C3 ⊆ C2, so "
     "C4 = C2. Fusion is still protective: it guarantees a recall floor "
     "when traces are sparse (C4 ≥ C3 always)."),
    ("Why C5 wins overall", AMBER,
     "Provenance-aware synthesis preserves C4's perfect recall while "
     "partially fixing static over-approximation on S3 (precision 0.60 → "
     "0.75) → best overall F1 = 0.964."),
    ("The one honest miss", RED,
     "C5's only imprecision: a real but out-of-scope flush call on S3. "
     "C2's element list omits the calling function, so 'suppressed-in-scope' "
     "and 'real-out-of-scope' look identical to the agent."),
]
positions = [(Inches(0.55), Inches(1.5)), (Inches(6.7), Inches(1.5)),
             (Inches(0.55), Inches(4.15)), (Inches(6.7), Inches(4.15))]
for (title, col, body), (px, py) in zip(cards, positions):
    panel(s, px, py, Inches(6.05), Inches(2.4))
    rect(s, px, py, Inches(0.12), Inches(2.4), fill=col)
    textbox(s, px + Inches(0.3), py + Inches(0.18), Inches(5.6), Inches(0.5),
            [{"text": title, "size": 15, "bold": True, "color": col}])
    textbox(s, px + Inches(0.3), py + Inches(0.7), Inches(5.6), Inches(1.6),
            [{"text": body, "size": 13.5, "color": DARK}])
footer(s, n)
note(s, """
[2:30] Four mechanistic points behind those numbers. Top-left: on S4 fusion is
additive because static misses the runtime connect() edge and dynamic misses the
LogLevel::NONE-suppressed edge — the union recovers both. Top-right: on S1–S3
it's redundant because trace edges are a subset of static edges — but still
protective as a recall floor. Bottom-left: C5 wins overall because
provenance-aware synthesis keeps C4's perfect recall while trimming static
over-approximation — on S3 it lifts precision from 0.60 to 0.75. Bottom-right,
being honest about the one miss: C5's only imprecision is a real but
out-of-scope flush call on S3. Because C2's element list records source and
target but not the calling function, "suppressed-but-in-scope" and
"real-but-out-of-scope" are indistinguishable to the agent. That's a concrete,
fixable limitation I'll return to.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 16 — H2
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "H2 — Guided Instrumentation Improves Focus",
                       kicker="Experiment · RQ2", num=n)
textbox(s, Inches(0.55), Inches(1.35), Inches(12), Inches(0.5),
        [{"text": "Full (probe everything, filter after) vs Guided (scope first, "
          "then probe only scenario-relevant call sites), on S1–S3.",
          "size": 15, "bold": True, "color": NAVY}])
rows = [
    ["Metric", "Full", "Guided", "Change"],
    ["Nodes", "6.0", "3.0", "−50%"],
    ["Edges", "9.0", "2.7", "−70%"],
    ["Coverage", "1.000", "1.000", "0% loss"],
    ["Noise", "0.704", "0.000", "−100%"],
]
table(s, Inches(0.55), Inches(2.05), Inches(6.0), rows, [1.9, 1.2, 1.2, 1.4],
      row_h=Inches(0.6), font_size=14, highlight_rows=[4])
panel(s, Inches(6.85), Inches(2.05), Inches(5.9), Inches(3.0), fill=LIGHT)
textbox(s, Inches(7.1), Inches(2.2), Inches(5.4), Inches(0.4),
        [{"text": "BLINDED 6-RATER EXPERT PANEL (mean, 1–5)", "size": 12,
          "bold": True, "color": BLUE}])
rows2 = [
    ["Dimension", "Full", "Guided"],
    ["Accuracy", "2.39", "4.61"],
    ["Readability", "2.44", "4.89"],
    ["Usefulness", "2.94", "4.00"],
]
table(s, Inches(7.1), Inches(2.7), Inches(5.4), rows2, [2.0, 1.0, 1.0],
      row_h=Inches(0.5), font_size=13)
panel(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.5),
      fill=RGBColor(0xEC,0xF6,0xEE))
textbox(s, Inches(0.8), Inches(5.5), Inches(11.9), Inches(1.3),
        [{"text": "Across all 54 rater × scenario × dimension comparisons: "
          "guided rated higher in 47, tied in 6, lower in exactly 1.",
          "size": 15, "bold": True, "color": GREEN, "bullet": True,
          "space_after": 4},
         {"text": "The single dissent (S2 usefulness): collapsing 7 repeated "
          "Config→Logging reads into one edge loses a count some find useful "
          "for debugging — an honest trade-off, not a defect.", "size": 13.5,
          "color": DARK, "bullet": True}])
footer(s, n)
note(s, """
[2:30] H2 answers RQ2 — focus. I compare "full" instrumentation, which probes
everything and filters afterwards, against "guided", which scopes first and only
probes the scenario-relevant call sites. The objective metrics on the left:
guided cuts nodes by 50% and edges by 70%, drives measured noise from 0.70 to
zero, and — critically — loses zero coverage. So it's not just smaller; nothing
relevant is dropped.
That's backed by a blinded six-rater expert panel, right. Guided wins on
accuracy, readability, and usefulness. Across all 54 comparisons, guided was
rated higher in 47, tied in 6, and lower in exactly one. That single dissent was
S2 usefulness, where collapsing seven repeated config-to-logging reads into one
edge loses a count some raters want for debugging — an honest trade-off, and I
report it as such.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 17 — H3
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "H3 — Is the Diagram Actually Useful?",
                       kicker="Experiment · Utility", num=n)
textbox(s, Inches(0.55), Inches(1.35), Inches(12), Inches(0.5),
        [{"text": "Does the pipeline's output (C5 interaction list) help an LLM "
          "agent comprehend a scenario — alone and in combination?",
          "size": 15, "bold": True, "color": NAVY}])
rows = [
    ["Condition", "Component recall", "Interaction recall", "Completeness (1–5)"],
    ["Source only", "0.917", "0.800", "4.25"],
    ["Trace only", "0.917", "0.900", "4.25"],
    ["Diagram only", "1.000", "1.000", "5.00"],
    ["Source + Trace", "0.917", "0.900", "4.25"],
    ["Source + Trace + Diagram", "1.000", "1.000", "5.00"],
]
table(s, Inches(0.55), Inches(2.0), Inches(12.2), rows, [3.0, 2.2, 2.2, 2.4],
      row_h=Inches(0.55), font_size=14, highlight_rows=[3, 5])
panel(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.2), fill=NAVY)
textbox(s, Inches(0.8), Inches(5.85), Inches(11.9), Inches(1.0),
        [{"text": "The standout: “Diagram only” equals “Source + Trace + Diagram” "
          "— perfect on every measure. With no source and no trace, the C5 "
          "diagram alone is a sufficient basis for scenario comprehension.",
          "size": 15, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, n)
note(s, """
[2:30] H3 asks the downstream question: is the reconstructed diagram actually
useful? I gave an LLM agent a scenario-comprehension task under five evidence
conditions and measured component recall, interaction recall, and a rubric
completeness score.
The standout result is in the two highlighted rows: "diagram only" and
"source + trace + diagram" are identical — and both are perfect. With no source
and no trace, just the C5 interaction list, the agent hits perfect component and
interaction recall and top completeness on every scenario, including the
constructed S4. Adding source and traces on top neither helps nor hurts, because
the diagram already saturates every metric. Notice source-only has the weakest
interaction recall at 0.80 — it foregrounds the main control-flow path and drops
secondary log-only branches, exactly the edges the fused diagram surfaces. So
the diagram isn't just a helpful supplement; it can be a *sufficient* basis for
comprehension.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 18 — Example S4 sequence
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Example Output — S4 Reconstructed Sequence",
                       kicker="Results · Qualitative", num=n)
IMG(s, os.path.join(FIG, "s4-sequence.png"),
    Inches(0.55), Inches(1.4), Inches(8.6), Inches(5.3))
panel(s, Inches(9.35), Inches(1.5), Inches(3.4), Inches(5.0), fill=LIGHT)
textbox(s, Inches(9.6), Inches(1.7), Inches(2.95), Inches(4.7),
        [{"text": "The dual-gap recovered", "size": 15, "bold": True,
          "color": BLUE, "space_after": 8},
         {"text": "stream — recovered from static evidence (runtime-suppressed "
          "by LogLevel::NONE).", "size": 13.5, "bullet": True,
          "space_after": 8},
         {"text": "onStageEvent — recovered from the trace (runtime connect(); "
          "no static call site).", "size": 13.5, "bullet": True,
          "space_after": 8},
         {"text": "Only the fused graph shows both edges in one scenario-"
          "specific SysML v2 diagram.", "size": 13.5, "bullet": True,
          "bold": True, "color": NAVY}])
footer(s, n)
note(s, """
[1:30] A concrete look at what the pipeline produces — the reconstructed S4
sequence. This makes the dual-gap tangible. The `stream` message is recovered
from static evidence, because at runtime it's suppressed by the LogLevel::NONE
no-op. The `onStageEvent` message is recovered from the trace, because it's
wired by a runtime connect() and has no static call site. Neither source alone
shows both — only the fused graph produces this single, scenario-specific
SysML v2 diagram with both edges present.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 19 — Industrial exploratory
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Exploratory: A Second, Larger Industrial Codebase",
                       kicker="Transferability", num=n)
panel(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.3), fill=LIGHT)
textbox(s, Inches(0.8), Inches(1.66), Inches(11.9), Inches(1.05),
        [{"text": "A practitioner session applied the same static-extraction and "
          "query approach to a second industrial system — testing whether the "
          "method transfers beyond CPSCore.", "size": 16, "color": DARK}],
        anchor=MSO_ANCHOR.MIDDLE)
panel(s, Inches(0.55), Inches(3.05), Inches(6.05), Inches(3.35),
      fill=RGBColor(0xEC,0xF6,0xEE))
textbox(s, Inches(0.85), Inches(3.25), Inches(5.5), Inches(3.0),
        [{"text": "Worked well", "size": 16, "bold": True, "color": GREEN,
          "space_after": 8},
         {"text": "Structural diagrams rated consistently useful.", "size": 15,
          "bullet": True, "space_after": 10},
         {"text": "Dependency diagrams rated consistently useful.", "size": 15,
          "bullet": True, "space_after": 10},
         {"text": "Static extraction + graph query transferred to a larger, "
          "unseen codebase.", "size": 15, "bullet": True, "bold": True,
          "color": NAVY}])
panel(s, Inches(6.7), Inches(3.05), Inches(6.05), Inches(3.35),
      fill=RGBColor(0xFB,0xF2,0xE3))
textbox(s, Inches(7.0), Inches(3.25), Inches(5.5), Inches(3.0),
        [{"text": "Clearest area to improve", "size": 16, "bold": True,
          "color": AMBER, "space_after": 8},
         {"text": "Sequence-diagram message labelling was identified as the "
          "weakest point.", "size": 15, "bullet": True, "space_after": 10},
         {"text": "Consistent with the CPSCore finding that message-level "
          "semantics need richer provenance.", "size": 15, "bullet": True,
          "space_after": 10},
         {"text": "A concrete, actionable direction for future work.",
          "size": 15, "bullet": True, "bold": True, "color": NAVY}])
footer(s, n)
note(s, """
[1:30] To probe transferability beyond CPSCore, a practitioner session applied
the same static-extraction and query approach to a second, larger industrial
codebase. The good news, on the left: the structural and dependency diagrams
were rated consistently useful — the static extraction and graph queries
transferred to an unseen, larger system. On the right, the clearest area to
improve was sequence-diagram message labelling, which is consistent with the
CPSCore finding that message-level semantics need richer provenance. It's
exploratory, not a controlled study, but it points at a concrete next step.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 20 — Unifying finding
# --------------------------------------------------------------------------- #
n = pg(); s = slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, Inches(0.3), SH, fill=AMBER)
textbox(s, Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.5),
        [{"text": "THE ONE PATTERN BEHIND ALL THREE RESULTS", "size": 15,
          "bold": True, "color": TEAL}])
textbox(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.6),
        [{"text": "Additional evidence improves reconstruction or comprehension "
          "specifically when the primary evidence source is missing something "
          "— and is otherwise redundant.", "size": 30, "bold": True,
          "color": WHITE}])
items = [
    ("H1", "Fusion is strictly additive only where static & dynamic each miss a "
           "different edge (S4); redundant where the static graph is complete."),
    ("H2", "Guiding instrumentation by the code graph removes 100% of noise "
           "with zero coverage loss."),
    ("H3", "The reconstructed diagram alone is a sufficient basis for scenario "
           "comprehension — perfect recall everywhere."),
]
y = Inches(4.5)
for tag, txt in items:
    chip(s, Inches(0.9), y, Inches(0.9), tag, AMBER, tcolor=NAVY,
         h=Inches(0.5), size=17)
    textbox(s, Inches(2.0), y - Inches(0.02), Inches(10.6), Inches(0.7),
            [{"text": txt, "size": 15, "color": RGBColor(0xDCE6F0)}],
            anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.78)
note(s, """
[1:00] Step back — the three results converge on a single pattern, and this is
the sentence I want you to remember: additional evidence improves reconstruction
or comprehension specifically when the primary source is missing something, and
is otherwise redundant. H1: fusion is additive only where the two sources each
miss a different edge, redundant where static is already complete. H2: guiding
instrumentation removes all noise with zero coverage loss. H3: the diagram alone
is a sufficient basis for comprehension. One principle, three angles.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 21 — Contributions
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Contributions", kicker="Summary", num=n)
contribs = [
    ("i", "One property graph holding static call graph + runtime trace "
          "together, remembering each interaction's provenance."),
    ("ii", "A scenario-scoping method, using plain scenario-scoped Cypher "
           "queries, that trims the combined graph to just the interactions "
           "relevant to one scenario."),
    ("iii", "A way to recover publish-subscribe wiring that ordinary "
            "call-graph analysis cannot see."),
    ("iv", "csp_matcher — a small, reversible tool that inserts and removes "
           "trace probes automatically."),
    ("v", "Eleven natural-language agent skills running the whole pipeline, "
          "source → SysML v2, on any codebase without code changes."),
    ("vi", "An evidence-based account of when combining static and dynamic "
           "analysis helps — and when it does not."),
]
y = Inches(1.5)
for tag, txt in contribs:
    panel(s, Inches(0.55), y, Inches(12.2), Inches(0.8), fill=LIGHT)
    chip(s, Inches(0.75), y + Inches(0.18), Inches(0.75), tag, BLUE,
         h=Inches(0.44), size=15)
    textbox(s, Inches(1.7), y + Inches(0.06), Inches(10.9), Inches(0.68),
            [{"text": txt, "size": 15, "color": DARK}], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.88)
footer(s, n)
note(s, """
[1:30] Six contributions. One: a single property graph that holds the static
call graph and the runtime trace together, remembering each interaction's
provenance. Two: a scenario-scoping method, built on plain scenario-scoped
Cypher queries, that trims that combined graph to just one scenario. Three: a
way to recover pub/sub wiring that call-graph analysis can't see. Four:
csp_matcher, the small reversible instrumentation tool. Five: eleven
natural-language agent skills that run the whole pipeline end-to-end, source to
SysML v2, on any codebase with no code changes. And six — which I'd argue is the
scientific contribution — a clear, evidence-based account of when combining
static and dynamic analysis actually helps, and when it doesn't.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 22 — Limitations & future work
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Limitations & Future Work",
                       kicker="Reflection", num=n)
panel(s, Inches(0.55), Inches(1.5), Inches(6.05), Inches(4.9),
      fill=RGBColor(0xFDECEA))
textbox(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(4.6),
        [{"text": "Limitations", "size": 17, "bold": True, "color": RED,
          "space_after": 8},
         {"text": "Single primary case study (CPSCore); S4 is constructed, not "
          "found in the wild.", "size": 14.5, "bullet": True, "space_after": 8},
         {"text": "C2's element list omits the calling function → agent cannot "
          "always distinguish in-scope from out-of-scope edges.", "size": 14.5,
          "bullet": True, "space_after": 8},
         {"text": "Two H3 conditions (source-only, diagram-only) rated non-blind "
          "by the same rater — disclosed.", "size": 14.5, "bullet": True,
          "space_after": 8},
         {"text": "Message-label semantics on sequence diagrams remain the "
          "weakest output dimension.", "size": 14.5, "bullet": True,
          "space_after": 8},
         {"text": "Traces limited to what the test suite exercises.", "size": 14.5,
          "bullet": True}])
panel(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.9),
      fill=RGBColor(0xEAF6F1))
textbox(s, Inches(7.15), Inches(1.7), Inches(5.35), Inches(4.6),
        [{"text": "Future work", "size": 17, "bold": True, "color": TEAL,
          "space_after": 8},
         {"text": "Carry the calling function into the element list to fix the "
          "in-scope / out-of-scope ambiguity.", "size": 14.5, "bullet": True,
          "space_after": 8},
         {"text": "Evaluate on systems with cross-module virtual dispatch, "
          "plugin loading, or reflection — where the additive regime dominates.",
          "size": 14.5, "bullet": True, "space_after": 8},
         {"text": "Richer, provenance-driven message labelling for sequence "
          "diagrams.", "size": 14.5, "bullet": True, "space_after": 8},
         {"text": "Broaden the expert panel and blind all conditions.",
          "size": 14.5, "bullet": True, "space_after": 8},
         {"text": "Continuous, CI-integrated model maintenance.", "size": 14.5,
          "bullet": True}])
footer(s, n)
note(s, """
[1:30] Being upfront about limitations, on the left. It's one primary case
study, and S4 is constructed. The in-scope/out-of-scope ambiguity comes from C2's
element list omitting the calling function. Two H3 conditions were rated
non-blind by the same rater, which I disclose. Message-label semantics are the
weakest output, and traces are bounded by what the test suite exercises.
Future work, right, maps almost one-to-one: carry the calling function into the
element list to remove that ambiguity; evaluate on systems with virtual
dispatch, plugins, or reflection, where the additive regime should dominate;
richer provenance-driven message labelling; broaden and fully blind the expert
panel; and integrate this into CI for continuous model maintenance.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 23 — Conclusion
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Conclusion", kicker="Wrap-up", num=n)
panel(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.2), fill=NAVY)
textbox(s, Inches(0.85), Inches(1.72), Inches(11.7), Inches(1.9),
        [{"text": "A hybrid pipeline fuses a static call graph with runtime "
          "traces in one provenance-tagged property graph and scopes it to "
          "scenario-specific SysML v2 interaction diagrams.", "size": 18,
          "bold": True, "color": WHITE, "space_after": 8},
         {"text": "Fusion helps exactly where one source is blind (C5 macro-F1 "
          "0.964); guided scoping removes 100% of noise at full coverage; the "
          "diagram alone suffices for comprehension.", "size": 16,
          "color": RGBColor(0xC9,0xD6,0xE2)}])
finals = [
    ("RQ1", "Fusion is additive where sources disagree (S4), redundant "
            "elsewhere — best macro-F1 0.964.", BLUE),
    ("RQ2", "Scenario-scoped Cypher querying yields focused, noise-free "
            "diagrams with zero coverage loss.", TEAL),
]
y = Inches(4.0)
for tag, txt, col in finals:
    chip(s, Inches(0.6), y, Inches(1.05), tag, col, h=Inches(0.55),
         size=16, tcolor=WHITE if col != AMBER else NAVY)
    textbox(s, Inches(1.85), y, Inches(10.8), Inches(0.6),
            [{"text": txt, "size": 16, "bold": True, "color": NAVY}],
            anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.78)
footer(s, n)
note(s, """
[1:00] To conclude: I built a hybrid pipeline that fuses a static call graph with
runtime traces in one provenance-tagged property graph and scopes it to
scenario-specific SysML v2 diagrams. Fusion helps exactly where one source is
blind — best overall F1 of 0.96; guided scoping removes all noise at full
coverage; and the reconstructed diagram alone suffices for comprehension.
Mapping back: RQ1 — fusion is additive where sources disagree, best F1 0.964.
RQ2 — scenario-scoped querying yields focused, noise-free diagrams with zero
coverage loss. Thank you — I'm happy to take questions.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 24 — Thank you
# --------------------------------------------------------------------------- #
n = pg(); s = slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.06), fill=AMBER)
logo_right(s, LOGO_TNO, SW - Inches(0.7), Inches(0.6), Inches(0.9), on_panel=True)
logo(s, LOGO_MSC, Inches(0.9), Inches(6.5), Inches(0.6))
textbox(s, Inches(0.9), Inches(2.2), Inches(11.6), Inches(1.3),
        [{"text": "Thank you — Questions?", "size": 46, "bold": True,
          "color": WHITE}])
textbox(s, Inches(0.9), Inches(3.75), Inches(11.6), Inches(1.0),
        [{"text": "Hybrid Static–Trace Analysis for Interaction Model Synthesis "
          "in Event-Driven Systems", "size": 18, "italic": True,
          "color": RGBColor(0xC9,0xD6,0xE2)}])
textbox(s, Inches(0.9), Inches(5.3), Inches(11.6), Inches(1.0),
        [{"text": "Prathik Anand Krishnan", "size": 22, "bold": True,
          "color": WHITE, "space_after": 4},
         {"text": "TNO-ESI  ·  MSc Software Engineering, University of Amsterdam",
          "size": 14, "color": TEAL}])
textbox(s, Inches(2.0), Inches(6.58), Inches(6), Inches(0.5),
        [{"text": "University of Amsterdam", "size": 12, "bold": True,
          "color": TEAL}])
note(s, """
Closing slide for the Q&A. Common questions to be ready for: (1) Why only one
case study — transferability is addressed by the exploratory second codebase.
(2) Isn't S4 contrived — yes, deliberately, to isolate the additive regime that
naturally-occurring virtual dispatch / plugins would trigger. (3) LLM
reproducibility — the reference set was authored before scoring; provenance
tagging keeps the synthesis auditable. (4) Scalability of scenario scoping —
it's a linear Cypher edge filter. (5) Threats to validity in the expert panel — six raters, A/B
randomised and blinded for the reused conditions.
""")

# --------------------------------------------------------------------------- #
#  SLIDE 25 — Backup
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Backup — S1 Structural & Sequence Views",
                       kicker="Appendix", num=n)
IMG(s, os.path.join(FIG, "s1-structural.png"),
    Inches(0.55), Inches(1.5), Inches(6.0), Inches(5.0))
IMG(s, os.path.join(FIG, "s1-sequence.png"),
    Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0))
textbox(s, Inches(0.55), Inches(6.55), Inches(6.0), Inches(0.4),
        [{"text": "S1 structural view", "size": 11, "italic": True,
          "color": GREY, "align": PP_ALIGN.CENTER}])
textbox(s, Inches(6.8), Inches(6.55), Inches(6.0), Inches(0.4),
        [{"text": "S1 reconstructed sequence", "size": 11, "italic": True,
          "color": GREY, "align": PP_ALIGN.CENTER}])
footer(s, n)
note(s, """
Backup slide for Q&A: the S1 structural and reconstructed sequence views, if the
committee wants to see a happy-path example end to end alongside the S4 fault
path shown earlier.
""")

# --------------------------------------------------------------------------- #
#  SLIDE — Backup: Engineering challenges  (from proposal slide 12)
# --------------------------------------------------------------------------- #
n = pg(); s = slide(); title_bar(s, "Backup — Engineering Challenges & Solutions",
                       kicker="Appendix", num=n)
rows = [
    ["Challenge", "How it was solved"],
    ["C++ is context-sensitive — regex / simple parsers fail",
     "Renaissance + Neo4j for semantic-level extraction; ClangSharp (libclang) for validation"],
    ["Thousands of dependencies crash SysML toolchains",
     "Identifier quoting + relationship deduplication in the generation layer"],
    ["Hardcoded scenario scoping is brittle",
     "Agent-skill architecture — scenario defined as a natural-language query"],
    ["Traces alone are incomplete",
     "Static reachability in the graph fills unobserved interactions — the core hybrid contribution"],
]
table(s, Inches(0.55), Inches(1.55), Inches(12.2), rows, [4.6, 7.4],
      row_h=Inches(1.0), font_size=14)
footer(s, n)
note(s, """
Backup slide for questions on engineering depth. Four concrete challenges and how
each was handled: C++ context-sensitivity forced semantic extraction via
Renaissance and Neo4j, with a ClangSharp/libclang cross-check; large dependency
counts that crash SysML tools were tamed with identifier quoting and relationship
deduplication; brittle hardcoded scoping was replaced by the agent-skill
architecture where a scenario is just a natural-language query; and trace
incompleteness is exactly what static reachability in the graph compensates for —
the core hybrid contribution.
""")

# --------------------------------------------------------------------------- #
#  Save
# --------------------------------------------------------------------------- #
out = os.path.join(HERE, "Thesis_Defence_Prathik.pptx")
try:
    prs.save(out)
except PermissionError:
    alt = os.path.join(HERE, "Thesis_Defence_Prathik_NEW.pptx")
    prs.save(alt)
    out = alt
    print("NOTE: target was locked (open in PowerPoint); saved to alternate file.")
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
